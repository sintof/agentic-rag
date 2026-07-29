"""
Document loading + chunking. Text extraction for txt/docx/pptx and the chunker are
adapted from the TezYodla project (D:\\planner\\projects\\tezyodla\\tfbk\\app\\document_parser.py) —
same document types, same proven chunking approach. PDF handling is new here: PyMuPDF
(fitz) pulls both text AND embedded images per page, so images get captioned by a
vision model and become searchable (the multimodal requirement).
"""
import base64
import re
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx"}


def is_supported_file(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def get_file_type(filename: str) -> str:
    return Path(filename).suffix.lower().lstrip(".")


def caption_image(image_bytes: bytes, mime: str = "image/png") -> str:
    """Caption a PDF-embedded image with the vision-capable lite model."""
    from .llm import get_llm_lite  # local import: avoids a hard dependency at module load

    b64 = base64.b64encode(image_bytes).decode("utf-8")
    llm = get_llm_lite(temperature=0.0)
    message = {
        "role": "user",
        "content": [
            {
                "type": "text",
                "text": (
                    "Describe this image/diagram in 1-3 sentences, focused on any facts, "
                    "labels, numbers, or relationships a student would need to answer a "
                    "question about it. Do not add information that isn't visible in the image."
                ),
            },
            {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
        ],
    }
    try:
        response = llm.invoke([message])
        return (response.content or "").strip()
    except Exception as exc:  # vision call can fail independently of the rest of ingestion
        return f"[Image present, captioning failed: {exc}]"


def extract_text_and_images(file_path: str) -> str:
    """PDF-specific: pulls page text AND embedded images (captioned) into one text stream."""
    doc = fitz.open(file_path)
    parts: list[str] = []
    for page_index, page in enumerate(doc, start=1):
        page_text = page.get_text().strip()
        if page_text:
            parts.append(page_text)

        for image_index, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception:
                continue
            image_bytes = base_image.get("image")
            ext = base_image.get("ext", "png")
            if not image_bytes:
                continue
            caption = caption_image(image_bytes, mime=f"image/{ext}")
            parts.append(f"[Image on page {page_index}]: {caption}")

    doc.close()
    return "\n\n".join(parts)


def extract_text(file_path: str, file_type: str) -> str:
    path = Path(file_path)
    if file_type == "txt":
        return path.read_text(encoding="utf-8", errors="ignore")
    if file_type == "pdf":
        return extract_text_and_images(str(path))
    if file_type == "docx":
        document = DocxDocument(str(path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    if file_type == "pptx":
        presentation = Presentation(str(path))
        slide_texts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = [f"Slide {index}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            slide_texts.append("\n".join(parts))
        return "\n\n".join(slide_texts)
    return ""


def chunk_text(text: str, max_words: int = 250, overlap_words: int = 40) -> list[dict]:
    """Paragraph-aware chunker with overlap. Smaller default chunk size than TezYodla's
    (250 vs 800 words) since these chunks feed a retriever with a top-K cutoff, not a
    whole-document evidence-extraction pass."""
    cleaned = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not cleaned:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", cleaned) if p.strip()]
    chunks: list[dict] = []
    current: list[str] = []
    current_words = 0

    for paragraph in paragraphs:
        words = paragraph.split()
        if current and current_words + len(words) > max_words:
            chunk_text_value = "\n\n".join(current)
            chunks.append({
                "chunk_id": f"chunk_{len(chunks) + 1:03d}",
                "text": chunk_text_value,
            })
            overlap = " ".join(chunk_text_value.split()[-overlap_words:]) if overlap_words else ""
            current = [overlap] if overlap else []
            current_words = len(overlap.split())
        current.append(paragraph)
        current_words += len(words)

    if current:
        chunks.append({
            "chunk_id": f"chunk_{len(chunks) + 1:03d}",
            "text": "\n\n".join(current),
        })

    return chunks
